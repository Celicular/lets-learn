import os
import pymupdf  # fitz
from pptx import Presentation

def extract_text_from_pdf(filepath):
    text = ""
    try:
        doc = pymupdf.open(filepath)
        for page in doc:
            text += page.get_text() + "\n\n"
        doc.close()
    except Exception as e:
        print(f"❌ Error parsing PDF {filepath}: {e}")
    print(f"📄 [PARSER] Extracted {len(text)} characters from PDF.")
    return text.strip()

def extract_text_from_pptx(filepath):
    text = ""
    try:
        prs = Presentation(filepath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            text += "\n"
    except Exception as e:
        print(f"❌ Error parsing PPTX {filepath}: {e}")
    print(f"📄 [PARSER] Extracted {len(text)} characters from PPTX.")
    return text.strip()

def extract_text_from_txt(filepath):
    text = ""
    try:
        with open(filepath, 'r', encoding='utf-8') as f:
            text = f.read()
    except Exception as e:
        print(f"❌ Error reading Text file {filepath}: {e}")
    print(f"📄 [PARSER] Loaded {len(text)} characters from {os.path.basename(filepath)}.")
    return text.strip()

def parse_document(filepath):
    """
    Routes the file to the appropriate parser based on the extension.
    Supported: .pdf, .pptx, .txt, .md
    """
    if not os.path.exists(filepath):
        print(f"Error: File not found at {filepath}")
        return ""
        
    ext = os.path.splitext(filepath)[1].lower()
    
    if ext == '.pdf':
        return extract_text_from_pdf(filepath)
    elif ext == '.pptx':
        return extract_text_from_pptx(filepath)
    elif ext in ['.txt', '.md', '.csv', '.json']:
        return extract_text_from_txt(filepath)
    else:
        print(f"Unsupported file type: {ext}")
        return ""
